from __future__ import annotations

from pathlib import Path
from typing import Any
import re

import pandas as pd
from docx import Document
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".csv", ".xlsx", ".xlsm", ".docx", ".pptx", ".txt", ".md"}
TEXT_COLUMN_HINTS = {
    "update_text",
    "update",
    "status",
    "comments",
    "notes",
    "description",
    "details",
    "activity",
    "task",
    "issue description",
    "problem tickets",
    "problem tickets / projects / crs",
    "project / cr #",
    "pr/cr/prb",
    "deliverables",
}
OWNER_COLUMN_HINTS = {"owner", "responsible", "assignee", "ba", "pm", "resource"}
WORKSTREAM_COLUMN_HINTS = {"workstream", "team", "phase", "market", "brand", "unit", "area / dept / team"}
PROJECT_COLUMN_HINTS = {"project", "project name", "project / cr #", "pr/cr/prb"}
STATUS_COLUMN_HINTS = {"status", "reporting status", "priority", "sign-off", "sign off", "% done", "progress"}


def normalize_text(text: Any) -> str:
    if text is None:
        return ""
    try:
        if pd.isna(text):
            return ""
    except (TypeError, ValueError):
        pass

    text = "" if text is None else str(text)
    text = re.sub(r"\s+", " ", text).strip()
    if text.lower() in {"nan", "nat", "none"}:
        return ""
    return text


def make_update_record(
    update_id: str,
    source_file: str,
    source_type: str,
    update_text: str,
    project: str = "Unknown Project",
    workstream: str = "General",
    owner: str = "Owner unclear",
    location: str = "",
) -> dict:
    return {
        "update_id": update_id,
        "source_file": source_file,
        "source_type": source_type,
        "source_location": location,
        "project": project or "Unknown Project",
        "workstream": workstream or "General",
        "update_text": normalize_text(update_text),
        "owner": owner or "Owner unclear",
    }


def parse_uploaded_file(uploaded_file) -> pd.DataFrame:
    filename = uploaded_file.name
    extension = Path(filename).suffix.lower()

    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {extension}")

    if extension == ".csv":
        return parse_csv(uploaded_file, filename)

    if extension in [".xlsx", ".xlsm"]:
        return parse_excel(uploaded_file, filename)

    if extension == ".docx":
        return parse_docx(uploaded_file, filename)

    if extension == ".pptx":
        return parse_pptx(uploaded_file, filename)

    if extension in [".txt", ".md"]:
        return parse_notes(uploaded_file, filename)

    raise ValueError(f"Unsupported file type: {extension}")


def parse_csv(uploaded_file, filename: str) -> pd.DataFrame:
    df = pd.read_csv(uploaded_file)
    return normalize_dataframe(df, filename, "csv")


def parse_excel(uploaded_file, filename: str) -> pd.DataFrame:
    excel_file = pd.ExcelFile(uploaded_file)
    records = []

    for sheet_name in excel_file.sheet_names:
        raw = excel_file.parse(sheet_name, header=None)
        normalized = normalize_excel_sheet(raw, filename, sheet_name=sheet_name)
        records.extend(normalized.to_dict("records"))

    return pd.DataFrame(records)


def normalize_header(value: Any) -> str:
    return normalize_text(value).lower().strip(":")


def cell_matches(header: str, hints: set[str]) -> bool:
    header = normalize_header(header)
    return header in hints or any(hint in header for hint in hints if len(hint) > 3)


def row_values(row: pd.Series) -> list[str]:
    return [normalize_text(value) for value in row.tolist() if normalize_text(value)]


def infer_project_name(raw: pd.DataFrame, filename: str, sheet_name: str = "") -> str:
    for _, row in raw.head(12).iterrows():
        values = [normalize_text(value) for value in row.tolist()]
        for idx, value in enumerate(values):
            if normalize_header(value) in {"project name", "project"}:
                for candidate in values[idx + 1 :]:
                    if candidate and not re.match(r"^(week starting|days)$", candidate, re.IGNORECASE):
                        return candidate

    name = Path(filename).stem
    lowered = name.lower()
    if "ebs" in lowered:
        return "EBS Development"
    if "ecomm" in lowered:
        return "InStore Ecomm Rollout"
    if "mpos" in lowered:
        return "MPOS Pilot Rollout"
    if "mdm" in lowered:
        if "product" in sheet_name.lower():
            return "Product MDM Process Mapping"
        if "location" in sheet_name.lower():
            return "Location MDM Process Mapping"
        return "MDM Process Mapping"
    return name


def detect_header_row(raw: pd.DataFrame) -> int | None:
    best_idx = None
    best_score = 0

    for idx, row in raw.head(25).iterrows():
        score = 0
        for value in row.tolist():
            header = normalize_header(value)
            if not header:
                continue
            if (
                cell_matches(header, TEXT_COLUMN_HINTS)
                or cell_matches(header, OWNER_COLUMN_HINTS)
                or cell_matches(header, WORKSTREAM_COLUMN_HINTS)
                or cell_matches(header, PROJECT_COLUMN_HINTS)
                or cell_matches(header, STATUS_COLUMN_HINTS)
                or header in {"start", "end", "eta", "date", "days"}
            ):
                score += 1

        if score > best_score:
            best_idx = idx
            best_score = score

    return best_idx if best_score >= 2 else None


def normalize_excel_sheet(raw: pd.DataFrame, filename: str, sheet_name: str = "") -> pd.DataFrame:
    records = []
    project_name = infer_project_name(raw, filename, sheet_name)
    header_idx = detect_header_row(raw)

    if header_idx is None:
        for idx, row in raw.iterrows():
            values = row_values(row)
            if not values:
                continue

            update_text = " | ".join(values)
            if should_skip_excel_update(update_text):
                continue

            records.append(
                make_update_record(
                    update_id=f"excel-{len(records) + 1}",
                    source_file=filename,
                    source_type="excel",
                    update_text=update_text,
                    project=project_name,
                    workstream=sheet_name or "General",
                    location=f"Sheet: {sheet_name}, Row: {idx + 1}",
                )
            )
        return pd.DataFrame(records)

    headers = [normalize_text(value) for value in raw.loc[header_idx].tolist()]
    data = raw.iloc[header_idx + 1 :].copy()

    for idx, row in data.iterrows():
        row_map = {
            headers[col_idx]: normalize_text(value)
            for col_idx, value in enumerate(row.tolist())
            if col_idx < len(headers) and headers[col_idx] and normalize_text(value)
        }
        values = [value for value in row_map.values() if value]
        if not values:
            continue

        update_text = build_excel_update_text(row_map, values)
        if should_skip_excel_update(update_text):
            continue

        records.append(
            make_update_record(
                update_id=f"excel-{len(records) + 1}",
                source_file=filename,
                source_type="excel",
                update_text=update_text,
                project=first_matching_value(row_map, PROJECT_COLUMN_HINTS) or project_name,
                workstream=first_matching_value(row_map, WORKSTREAM_COLUMN_HINTS) or sheet_name or "General",
                owner=first_matching_value(row_map, OWNER_COLUMN_HINTS) or "Owner unclear",
                location=f"Sheet: {sheet_name}, Row: {idx + 1}",
            )
        )

    return pd.DataFrame(records)


def build_excel_update_text(row_map: dict[str, str], values: list[str]) -> str:
    parts = []
    for header, value in row_map.items():
        if cell_matches(header, TEXT_COLUMN_HINTS) or cell_matches(header, STATUS_COLUMN_HINTS):
            parts.append(f"{header}: {value}")

    if not parts:
        for header, value in row_map.items():
            if not looks_like_timeline_marker(value):
                parts.append(f"{header}: {value}" if header else value)

    return "; ".join(parts) if parts else " | ".join(values)


def first_matching_value(row_map: dict[str, str], hints: set[str]) -> str:
    for header, value in row_map.items():
        if cell_matches(header, hints) and value:
            return value
    return ""


def looks_like_timeline_marker(value: str) -> bool:
    if value.lower() in {"u", "p", "x", "-", "yes", "no"}:
        return True
    if re.fullmatch(r"\d+(\.0)?", value):
        return True
    return False


def should_skip_excel_update(update_text: str) -> bool:
    text = normalize_text(update_text)
    if len(text) < 8:
        return True
    values = [part.strip() for part in re.split(r"[|;]", text) if part.strip()]
    if not values:
        return True
    timeline_markers = sum(1 for value in values if looks_like_timeline_marker(value))
    return timeline_markers >= max(3, len(values) - 1)


def normalize_dataframe(
    df: pd.DataFrame,
    filename: str,
    source_type: str,
    sheet_name: str = "",
) -> pd.DataFrame:
    records = []

    lower_cols = {str(col).lower().strip(): col for col in df.columns}

    text_col = next((lower_cols[col] for col in lower_cols if cell_matches(col, TEXT_COLUMN_HINTS)), None)

    project_col = next((lower_cols[col] for col in lower_cols if cell_matches(col, PROJECT_COLUMN_HINTS)), None)
    workstream_col = next((lower_cols[col] for col in lower_cols if cell_matches(col, WORKSTREAM_COLUMN_HINTS)), None)
    owner_col = next((lower_cols[col] for col in lower_cols if cell_matches(col, OWNER_COLUMN_HINTS)), None)

    for idx, row in df.iterrows():
        if text_col:
            update_text = normalize_text(row.get(text_col, ""))
            if not update_text:
                values = [normalize_text(v) for v in row.values if normalize_text(v)]
                update_text = " | ".join(values)
        else:
            values = [normalize_text(v) for v in row.values if normalize_text(v)]
            update_text = " | ".join(values)

        if not update_text:
            continue

        project = normalize_text(row.get(project_col, "")) if project_col else "Unknown Project"
        workstream = normalize_text(row.get(workstream_col, "")) if workstream_col else "General"
        owner = normalize_text(row.get(owner_col, "")) if owner_col else "Owner unclear"

        location = f"Sheet: {sheet_name}, Row: {idx + 2}" if sheet_name else f"Row: {idx + 2}"

        records.append(
            make_update_record(
                update_id=f"{source_type}-{len(records) + 1}",
                source_file=filename,
                source_type=source_type,
                update_text=update_text,
                project=project,
                workstream=workstream,
                owner=owner,
                location=location,
            )
        )

    return pd.DataFrame(records)


def parse_docx(uploaded_file, filename: str) -> pd.DataFrame:
    document = Document(uploaded_file)
    records = []

    for idx, paragraph in enumerate(document.paragraphs, start=1):
        text = normalize_text(paragraph.text)
        if len(text) < 10:
            continue

        records.append(
            make_update_record(
                update_id=f"docx-p{idx}",
                source_file=filename,
                source_type="docx",
                update_text=text,
                location=f"Paragraph {idx}",
            )
        )

    table_counter = 0
    for table in document.tables:
        table_counter += 1
        for row_idx, row in enumerate(table.rows, start=1):
            cells = [normalize_text(cell.text) for cell in row.cells if normalize_text(cell.text)]
            text = " | ".join(cells)
            if len(text) < 10:
                continue

            records.append(
                make_update_record(
                    update_id=f"docx-t{table_counter}-r{row_idx}",
                    source_file=filename,
                    source_type="docx_table",
                    update_text=text,
                    location=f"Table {table_counter}, Row {row_idx}",
                )
            )

    return pd.DataFrame(records)


def parse_pptx(uploaded_file, filename: str) -> pd.DataFrame:
    presentation = Presentation(uploaded_file)
    records = []

    for slide_idx, slide in enumerate(presentation.slides, start=1):
        slide_text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = normalize_text(shape.text)
                if text:
                    slide_text_parts.append(text)

        slide_text = " | ".join(slide_text_parts)

        if len(slide_text) < 10:
            continue

        records.append(
            make_update_record(
                update_id=f"pptx-s{slide_idx}",
                source_file=filename,
                source_type="pptx",
                update_text=slide_text,
                location=f"Slide {slide_idx}",
            )
        )

    return pd.DataFrame(records)


def parse_notes(uploaded_file, filename: str) -> pd.DataFrame:
    raw = uploaded_file.read()

    if isinstance(raw, bytes):
        text = raw.decode("utf-8", errors="ignore")
    else:
        text = str(raw)

    chunks = chunk_notes_text(text)
    records = []

    for idx, chunk in enumerate(chunks, start=1):
        records.append(
            make_update_record(
                update_id=f"note-{idx}",
                source_file=filename,
                source_type="notes",
                update_text=chunk,
                location=f"Chunk {idx}",
            )
        )

    return pd.DataFrame(records)


def chunk_notes_text(text: str) -> list[str]:
    lines = [normalize_text(line) for line in text.splitlines()]
    lines = [line for line in lines if line]

    chunks = []
    buffer = []

    for line in lines:
        starts_new_item = line.startswith("- ") or line.startswith("* ") or re.match(r"^\d+[\.\)]\s+", line)

        if starts_new_item and buffer:
            chunks.append(" ".join(buffer))
            buffer = [line]
        else:
            buffer.append(line)

    if buffer:
        chunks.append(" ".join(buffer))

    final_chunks = []
    for chunk in chunks:
        if len(chunk) > 800:
            final_chunks.extend(split_long_text(chunk, max_chars=800))
        else:
            final_chunks.append(chunk)

    return [chunk for chunk in final_chunks if len(chunk) >= 10]


def split_long_text(text: str, max_chars: int = 800) -> list[str]:
    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks = []
    current = ""

    for sentence in sentences:
        if len(current) + len(sentence) <= max_chars:
            current = f"{current} {sentence}".strip()
        else:
            if current:
                chunks.append(current)
            current = sentence

    if current:
        chunks.append(current)

    return chunks
